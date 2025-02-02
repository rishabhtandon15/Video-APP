import os
from flask import Flask, request, render_template, send_file
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from transformers import AutoTokenizer, AutoModelForCausalLM
import torch

app = Flask(__name__)

# Load the Qwen model and tokenizer
tokenizer = AutoTokenizer.from_pretrained("Qwen/Qwen2.5-1.5B-Instruct")
model = AutoModelForCausalLM.from_pretrained("Qwen/Qwen2.5-1.5B-Instruct")

def generate_content(prompt, max_tokens=800):
    """Generates structured content for a PowerPoint presentation from the given prompt."""
    
    # Encode the prompt
    inputs = tokenizer(prompt, return_tensors="pt", padding=True, truncation=True)
    
    # Generate text from the model
    with torch.no_grad():
        outputs = model.generate(
            inputs['input_ids'], 
            attention_mask=inputs['attention_mask'], 
            max_length=max_tokens, 
            num_return_sequences=1
        )
    
    # Decode the generated text
    response = tokenizer.decode(outputs[0], skip_special_tokens=True)
    return response.strip()  # Remove unnecessary leading/trailing spaces

def create_ppt_from_content(content):
    """Creates a PowerPoint presentation from the structured content."""
    presentation = Presentation()
    
    # Split content into slides based on new lines or paragraphs
    slides = content.split("\n\n")  # Split by two newlines (indicating different sections)

    for slide_text in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
        parts = slide_text.split(":")  # Expecting "Title: Content"
        
        if len(parts) >= 2:
            title, body = parts[0], ":".join(parts[1:])
        else:
            title, body = parts[0], ""  # Handle cases where no colon exists

        # Set title formatting
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title.strip()
            title_text_frame = title_shape.text_frame
            title_text_frame.paragraphs[0].font.size = Pt(22)  
            title_text_frame.paragraphs[0].font.bold = True  
            title_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  

        # Set body text formatting and prevent overlap
        body_shape = slide.shapes.placeholders[1]
        if body_shape:
            body_shape.text = body.strip()
            body_text_frame = body_shape.text_frame
            # Dynamically adjust font size based on length of content
            for paragraph in body_text_frame.paragraphs:
                paragraph.font.size = Pt(18)  # Default font size
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  
            
            # Resize the body text box based on content length
            text_length = len(body.strip().split())
            if text_length > 300:
                body_shape.width = Inches(8.0)  # Adjust width of the text box if content is long
                body_shape.height = Inches(4.5)  # Adjust height of the text box
                body_text_frame.word_wrap = True
            elif text_length > 150:
                body_shape.width = Inches(8.0)  # Adjust width
                body_shape.height = Inches(3.5)  # Adjust height
                body_text_frame.word_wrap = True
            else:
                body_shape.width = Inches(8.0)
                body_shape.height = Inches(3.0)

    # Save the PowerPoint file temporarily
    ppt_filename = "generated_presentation.pptx"
    presentation.save(ppt_filename)
    
    return ppt_filename

@app.route("/", methods=["GET"])
def index():
    return render_template("3.html")  # Serve the HTML page

@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    # Get the content and theme from the form data
    content = request.form["content"]
    theme = request.form["theme"]
    
    # Generate content based on the theme
    prompt = f"Create a PowerPoint presentation on the following topic: {content}. The theme should be {theme}."
    generated_content = generate_content(prompt)
    
    # Create the PowerPoint presentation from the generated content
    ppt_filename = create_ppt_from_content(generated_content)
    
    # Serve the file as a downloadable response with a "Save As" dialog
    return send_file(ppt_filename, as_attachment=True, download_name="generated_presentation.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

if __name__ == "__main__":
    app.run(debug=True)
